from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT / "src"))


BLUE = "145DA0"
DARK = "1F2937"
GREEN = "2E7D32"
GRAY = "6B7280"


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_metrics(path: Path) -> dict[str, object]:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = rgb(DARK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.95), Inches(12.0), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(10)
        sp.font.color.rgb = rgb(GRAY)


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(DARK)
        p.space_after = Pt(7)


def add_metric_card(slide, label: str, value: str, left: float, top: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.25), Inches(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("EAF4FF")
    shape.line.color.rgb = rgb("C7DDF5")
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = rgb(BLUE)
    p.alignment = PP_ALIGN.CENTER
    q = frame.add_paragraph()
    q.text = label
    q.font.size = Pt(9)
    q.font.color.rgb = rgb(GRAY)
    q.alignment = PP_ALIGN.CENTER


def add_image_if_exists(slide, path: Path, left: float, top: float, width: float) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def make_deck(metrics: dict[str, object], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(BLUE)
        bar.line.fill.background()
        return s

    s = slide()
    add_title(s, "ChurnZero 26: Banking Customer Churn Prediction", "Round 2 case-study submission | IIT Kharagpur")
    add_bullets(
        s,
        [
            "Objective: predict customers likely to leave the bank and prioritize retention action.",
            "Deliverables: reproducible GitHub code, test-set predictions CSV, and presentation.",
            "Solution: explainable tabular ML pipeline with validation metrics and risk segmentation.",
        ],
        0.9,
        2.0,
        8.8,
        2.3,
        22,
    )
    add_metric_card(s, "Validation ROC-AUC", f"{metrics.get('roc_auc', 0):.3f}" if metrics else "TBD", 10.1, 1.55)
    add_metric_card(s, "Avg Precision", f"{metrics.get('average_precision', 0):.3f}" if metrics else "TBD", 10.1, 2.7)
    add_metric_card(s, "F1", f"{metrics.get('f1', 0):.3f}" if metrics else "TBD", 10.1, 3.85)

    s = slide()
    add_title(s, "Problem Framing", "Binary classification with business ranking output")
    add_bullets(
        s,
        [
            "Churn is modeled as the positive class: 1 means customer is likely to exit.",
            "The model outputs probabilities for customer-level prioritization, not only hard labels.",
            "Business value comes from finding high-risk customers early enough for targeted retention.",
            "Evaluation focuses on ranking quality, recall on churners, and stable validation performance.",
        ],
        0.9,
        1.55,
        11.5,
        4.8,
    )

    s = slide()
    add_title(s, "Data Understanding", "Banking profile, activity, and product usage signals")
    summary = metrics.get("data_summary", {}) if metrics else {}
    add_metric_card(s, "Rows", str(summary.get("rows", "TBD")), 0.9, 1.4)
    add_metric_card(s, "Columns", str(summary.get("columns", "TBD")), 3.4, 1.4)
    add_metric_card(s, "Target Rate", f"{100 * metrics.get('target_rate_train', 0):.1f}%" if metrics else "TBD", 5.9, 1.4)
    add_metric_card(s, "Features Used", str(metrics.get("feature_count", "TBD")) if metrics else "TBD", 8.4, 1.4)
    add_bullets(
        s,
        [
            "Numerical fields: age, tenure, balance, credit score, salary, product counts, activity flags.",
            "Categorical fields: geography, gender, product/customer segment fields if present.",
            "Identifiers and high-cardinality name-like fields are excluded to reduce leakage and overfit.",
        ],
        0.95,
        3.0,
        11.2,
        2.6,
    )

    s = slide()
    add_title(s, "Preprocessing Pipeline", "Reproducible transformations inside a single sklearn pipeline")
    add_bullets(
        s,
        [
            "Median imputation for numeric features.",
            "Most-frequent imputation and one-hot encoding for categorical features.",
            "Unknown categories are ignored safely during inference.",
            "The same fitted pipeline is saved and reused for official test predictions.",
        ],
        1.0,
        1.45,
        11.0,
        4.8,
    )

    s = slide()
    add_title(s, "Model Strategy", "Soft-voting ensemble tuned for tabular churn data")
    add_bullets(
        s,
        [
            "Extra Trees captures nonlinear interactions and robust split behavior.",
            "Histogram Gradient Boosting improves probability ranking on dense tabular patterns.",
            "Class-weighting handles churn imbalance without duplicating rows.",
            "Validation threshold is selected to maximize F1 while keeping probabilities for ranking.",
        ],
        0.95,
        1.45,
        11.4,
        4.8,
    )

    s = slide()
    add_title(s, "Validation Performance", "Holdout split with stratification")
    add_metric_card(s, "ROC-AUC", f"{metrics.get('roc_auc', 0):.3f}" if metrics else "TBD", 0.9, 1.45)
    add_metric_card(s, "Avg Precision", f"{metrics.get('average_precision', 0):.3f}" if metrics else "TBD", 3.4, 1.45)
    add_metric_card(s, "Precision", f"{metrics.get('precision', 0):.3f}" if metrics else "TBD", 5.9, 1.45)
    add_metric_card(s, "Recall", f"{metrics.get('recall', 0):.3f}" if metrics else "TBD", 8.4, 1.45)
    add_metric_card(s, "F1", f"{metrics.get('f1', 0):.3f}" if metrics else "TBD", 10.9, 1.45)
    add_image_if_exists(s, PROJECT_ROOT / "reports/figures/roc_curve.png", 0.95, 3.0, 5.6)
    add_image_if_exists(s, PROJECT_ROOT / "reports/figures/precision_recall_curve.png", 6.8, 3.0, 5.6)

    s = slide()
    add_title(s, "Error Analysis", "Confusion matrix at tuned decision threshold")
    add_image_if_exists(s, PROJECT_ROOT / "reports/figures/confusion_matrix.png", 0.95, 1.35, 6.3)
    add_bullets(
        s,
        [
            "False negatives represent missed retention opportunities.",
            "False positives are lower-cost if handled through low-friction engagement.",
            "Threshold can be adjusted depending on campaign budget and retention capacity.",
        ],
        7.6,
        1.7,
        5.0,
        3.5,
        16,
    )

    s = slide()
    add_title(s, "Churn Drivers", "Permutation importance on validation split")
    add_image_if_exists(s, PROJECT_ROOT / "reports/figures/feature_importance.png", 0.9, 1.25, 7.0)
    top_features = [item["feature"] for item in metrics.get("top_features", [])[:6]] if metrics else []
    add_bullets(
        s,
        [f"Top signal: {feature}" for feature in top_features] or ["Feature drivers will populate after training."],
        8.2,
        1.55,
        4.5,
        4.8,
        14,
    )

    s = slide()
    add_title(s, "Retention Playbook", "Convert risk scores into business action")
    add_bullets(
        s,
        [
            "Top 5% risk: proactive relationship-manager outreach and complaint-resolution review.",
            "5-20% risk: personalized offers based on product usage, tenure, and balance behavior.",
            "20-40% risk: education nudges, digital engagement, and service-quality checks.",
            "Low risk: monitor only; avoid unnecessary incentive spend.",
        ],
        0.9,
        1.45,
        11.6,
        4.8,
    )

    s = slide()
    add_title(s, "Operational Deployment", "Batch scoring flow for official test and future monthly refreshes")
    add_bullets(
        s,
        [
            "Ingest customer snapshot CSV.",
            "Run saved preprocessing and ensemble model.",
            "Export customer ID, churn probability, and churn prediction.",
            "Review top-risk segments and launch retention campaigns.",
            "Retrain monthly or whenever drift is detected in score distributions.",
        ],
        0.95,
        1.45,
        11.2,
        4.8,
    )

    s = slide()
    add_title(s, "Reproducibility", "Repository commands used for the Round 2 submission")
    add_bullets(
        s,
        [
            "python scripts/make_sample_data.py",
            "python scripts/train.py --train data/raw/train.csv --test data/raw/test.csv --target Exited",
            "python scripts/predict.py --model models/churn_model.joblib --test data/raw/test.csv --out outputs/predictions.csv",
            "python scripts/make_deck.py",
        ],
        0.95,
        1.55,
        11.7,
        4.5,
        15,
    )

    s = slide()
    add_title(s, "Risks and Controls", "What was guarded before final submission")
    add_bullets(
        s,
        [
            "Leakage: customer IDs, names, and high-cardinality text are excluded from features.",
            "Imbalance: class weighting and average precision are used alongside ROC-AUC.",
            "Generalization: stratified holdout validation and permutation importance avoid train-only claims.",
            "Inference safety: missing official-test columns are filled, unknown categories are handled.",
        ],
        0.95,
        1.45,
        11.5,
        4.8,
    )

    s = slide()
    add_title(s, "Final Submission Contents", "Round 2 checklist")
    add_bullets(
        s,
        [
            "PPTX/PDF deck: submission/ChurnZero_26_Round2_Presentation.pptx",
            "GitHub repository: code, README, requirements, scripts, tests.",
            "Predictions CSV: outputs/predictions.csv",
            "Model artifact: models/churn_model.joblib",
            "Validation evidence: outputs/metrics.json and reports/figures/*.png",
        ],
        0.95,
        1.45,
        11.6,
        4.8,
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Create Round 2 PowerPoint presentation.")
    parser.add_argument("--metrics", default="outputs/metrics.json")
    parser.add_argument("--out", default="submission/ChurnZero_26_Round2_Presentation.pptx")
    args = parser.parse_args()

    metrics = load_metrics(Path(args.metrics))
    out_path = Path(args.out)
    make_deck(metrics, out_path)
    print(f"Wrote deck: {out_path}")


if __name__ == "__main__":
    main()
